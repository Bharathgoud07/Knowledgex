# resources/views.py
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib.admin.views.decorators import staff_member_required
from django.contrib import messages
from django.db.models import Q, F, Count, Sum, Avg
from django.core.paginator import Paginator
from django.utils import timezone
from django.contrib.auth.models import User
from django.db.models.functions import TruncDate, Coalesce
import json
import zipfile

from .models import (
    Resource,
    SEMESTER_CHOICES,
    RESOURCE_TYPE_CHOICES,
    Subject,
    Favorite,
    Comment,
    Rating,
    Notification,
    Report,
    Visit,
)

from .forms import ResourceForm, CommentForm, RatingForm, ReportForm

# extra optional preview libs
try:
    import docx  # for DOCX preview
except ImportError:
    docx = None

try:
    import pptx  # for PPT/PPTX preview
except ImportError:
    pptx = None


# -------------------------------------------------------------------
# Resource list + filters
# -------------------------------------------------------------------
@login_required
def resource_list(request):
    resources = Resource.objects.all().order_by("-created_at")

    q = request.GET.get("q", "").strip()
    subject = request.GET.get("subject", "").strip()
    semester = request.GET.get("semester", "").strip()
    resource_type = request.GET.get("resource_type", "").strip()
    sort = request.GET.get("sort", "newest")

    # Search
    if q:
        resources = resources.filter(
            Q(title__icontains=q)
            | Q(description__icontains=q)
            | Q(subject__name__icontains=q)
            | Q(owner__username__icontains=q)
        )

    # Subject filter
    if subject:
        try:
            subject_id = int(subject)
            resources = resources.filter(subject_id=subject_id)
        except ValueError:
            pass

    # Semester filter
    if semester:
        try:
            sem_value = int(semester)
            resources = resources.filter(semester=sem_value)
        except ValueError:
            pass

    # Resource type filter
    if resource_type:
        resources = resources.filter(resource_type=resource_type)

    # Sorting
    if sort == "oldest":
        resources = resources.order_by("created_at")
    elif sort == "downloads":
        resources = resources.order_by("-download_count", "-created_at")
    elif sort == "az":
        resources = resources.order_by("title")
    elif sort == "subject":
        resources = resources.order_by("subject__name", "title")
    else:  # newest
        resources = resources.order_by("-created_at")

    paginator = Paginator(resources, 10)
    page_number = request.GET.get("page")
    page_obj = paginator.get_page(page_number)

    top_downloads = Resource.objects.order_by("-download_count").first()

    context = {
        "page_obj": page_obj,
        "q": q,
        "subject_filter": subject,
        "semester_filter": semester,
        "resource_type_filter": resource_type,
        "sort": sort,
        "semester_choices": SEMESTER_CHOICES,
        "resource_type_choices": RESOURCE_TYPE_CHOICES,
        "subjects": Subject.objects.order_by("name"),
        "top_downloads": top_downloads,
    }
    return render(request, "resources/resource_list.html", context)


# -------------------------------------------------------------------
# Upload resource
# -------------------------------------------------------------------
@login_required
def upload_resource(request):
    if request.method == "POST":
        form = ResourceForm(request.POST, request.FILES)
        if form.is_valid():
            resource = form.save(commit=False)
            resource.owner = request.user
            resource.save()
            messages.success(request, "Resource uploaded successfully!")
            return redirect("resource_list")
        else:
            messages.error(request, "Please correct the errors in the form.")
    else:
        form = ResourceForm()

    return render(request, "resources/upload_resource.html", {"form": form})


# -------------------------------------------------------------------
# Resource detail (comments + ratings + favorite)
# -------------------------------------------------------------------
@login_required
def resource_detail(request, pk):
    resource = get_object_or_404(Resource, pk=pk)

    # Count views only on GET
    if request.method == "GET":
        Resource.objects.filter(pk=pk).update(view_count=F("view_count") + 1)
        resource.refresh_from_db()

    comments = (
        resource.comments.select_related("user")
        .prefetch_related("replies__user")
        .filter(parent__isnull=True)
    )

    if request.method == "POST":
        # ----- New comment / reply -----
        if "comment_submit" in request.POST:
            comment_form = CommentForm(request.POST)
            rating_form = RatingForm()  # untouched
            if comment_form.is_valid():
                parent = None
                parent_id = request.POST.get("parent_id")
                if parent_id:
                    parent = Comment.objects.filter(
                        pk=parent_id, resource=resource
                    ).first()

                comment = comment_form.save(commit=False)
                comment.resource = resource
                comment.user = request.user
                comment.parent = parent
                comment.save()

                # Notify owners
                if resource.owner != request.user:
                    Notification.objects.create(
                        user=resource.owner,
                        notif_type="COMMENT" if parent is None else "REPLY",
                        message=(
                            f"{request.user.username} "
                            f"{'commented on' if parent is None else 'replied on'} "
                            f"your resource '{resource.title}'."
                        ),
                        resource=resource,
                        comment=comment,
                    )

                # Notify parent commenter if different
                if parent and parent.user not in (request.user, resource.owner):
                    Notification.objects.create(
                        user=parent.user,
                        notif_type="REPLY",
                        message=(
                            f"{request.user.username} replied to your comment "
                            f"on '{resource.title}'."
                        ),
                        resource=resource,
                        comment=comment,
                    )

                messages.success(request, "Comment added!")
                return redirect("resource_detail", pk=resource.pk)

        # ----- Rating submit -----
        elif "rating_submit" in request.POST:
            rating_form = RatingForm(request.POST)
            comment_form = CommentForm()
            if rating_form.is_valid():
                stars = rating_form.cleaned_data["stars"]
                Rating.objects.update_or_create(
                    resource=resource,
                    user=request.user,
                    defaults={"stars": stars},
                )

                if resource.owner != request.user:
                    Notification.objects.create(
                        user=resource.owner,
                        notif_type="RATING",
                        message=(
                            f"{request.user.username} rated your resource "
                            f"'{resource.title}' with {stars}★."
                        ),
                        resource=resource,
                    )

                messages.success(request, "Your rating has been saved!")
                return redirect("resource_detail", pk=resource.pk)
    else:
        comment_form = CommentForm()
        try:
            existing = Rating.objects.get(resource=resource, user=request.user)
            rating_form = RatingForm(initial={"stars": existing.stars})
        except Rating.DoesNotExist:
            rating_form = RatingForm()

    is_favorite = resource.is_favorited_by(request.user)

    context = {
        "resource": resource,
        "resource_type_choices": RESOURCE_TYPE_CHOICES,
        "comment_form": comment_form,
        "rating_form": rating_form,
        "comments": comments,
        "is_favorite": is_favorite,
    }
    return render(request, "resources/resource_detail.html", context)


# -------------------------------------------------------------------
# Download
# -------------------------------------------------------------------
@login_required
def resource_download(request, pk):
    resource = get_object_or_404(Resource, pk=pk)
    Resource.objects.filter(pk=pk).update(download_count=F("download_count") + 1)
    return redirect(resource.file.url)


# -------------------------------------------------------------------
# ONLINE VIEWER (image / pdf / docx / ppt / zip)
# -------------------------------------------------------------------
@login_required
def resource_viewer(request, pk):
    """
    Online viewer for various file types.
    """
    resource = get_object_or_404(Resource, pk=pk)

    # Count as a view
    Resource.objects.filter(pk=pk).update(view_count=F("view_count") + 1)
    resource.refresh_from_db()

    ext = (resource.file_ext or "").lower()
    preview_type = "fallback"
    preview_data = None

    try:
        # ---------- 1) IMAGE ----------
        if resource.is_image:
            preview_type = "image"
            preview_data = resource.file.url

        # ---------- 2) PDF ----------
        elif ext == "pdf":
            preview_type = "pdf"
            preview_data = resource.file.url

        # ---------- 3) DOCX ----------
        elif ext == "docx":
            preview_type = "docx"
            if docx is None:
                preview_data = None
            else:
                try:
                    document = docx.Document(resource.file.path)
                    lines = [p.text.strip() for p in document.paragraphs if p.text.strip()]
                    preview_data = lines[:40]  # first 40 lines
                except Exception:
                    preview_data = None

        # ---------- 4) PPT / PPTX ----------
        elif ext in ["ppt", "pptx"]:
            preview_type = "ppt"
            if pptx is None:
                preview_data = None
            else:
                try:
                    prs = pptx.Presentation(resource.file.path)
                    titles = []
                    for i, slide in enumerate(prs.slides, start=1):
                        text = ""
                        if slide.shapes.title and slide.shapes.title.text:
                            text = slide.shapes.title.text.strip()
                        else:
                            text = f"Slide {i}"
                        titles.append(text)
                    preview_data = titles
                except Exception:
                    preview_data = None

        # ---------- 5) ZIP ----------
        elif ext == "zip":
            preview_type = "zip"
            try:
                names = []
                with zipfile.ZipFile(resource.file.path, "r") as zf:
                    for info in zf.infolist():
                        if not info.is_dir():
                            names.append(info.filename)
                preview_data = names
            except Exception:
                preview_data = None

        # ---------- 6) OTHER FILES ----------
        else:
            preview_type = "fallback"
            preview_data = None

    except Exception:
        preview_type = "fallback"
        preview_data = None

    return render(
        request,
        "resources/resource_viewer.html",
        {
            "resource": resource,
            "preview_type": preview_type,
            "preview_data": preview_data,
        },
    )


# -------------------------------------------------------------------
# Favorites
# -------------------------------------------------------------------
@login_required
def toggle_favorite(request, pk):
    resource = get_object_or_404(Resource, pk=pk)
    favorite, created = Favorite.objects.get_or_create(
        user=request.user, resource=resource
    )
    if created:
        messages.success(request, "Added to your favorites!")
    else:
        favorite.delete()
        messages.info(request, "Removed from your favorites.")
    return redirect("resource_detail", pk=pk)


@login_required
def my_favorites(request):
    favorites = Favorite.objects.filter(user=request.user).select_related(
        "resource", "resource__subject"
    )
    resources = [f.resource for f in favorites]
    return render(
        request,
        "resources/my_favorites.html",
        {
            "resources": resources,
            "resource_type_choices": RESOURCE_TYPE_CHOICES,
        },
    )

# -------------------------------------------------------------------
# Subject dashboard (for charts)
# -------------------------------------------------------------------
@login_required
def subject_dashboard(request):
    subject_stats = (
        Resource.objects.values("subject__name")
        .annotate(
            total_uploads=Count("id"),
            total_downloads=Sum("download_count"),
            total_views=Sum("view_count"),
        )
        .order_by("subject__name")
    )

    labels = []
    uploads = []
    downloads = []
    views = []

    for row in subject_stats:
        labels.append(row["subject__name"] or "N/A")
        uploads.append(row["total_uploads"] or 0)
        downloads.append(row["total_downloads"] or 0)
        views.append(row["total_views"] or 0)

    context = {
        "labels": labels,
        "uploads": uploads,
        "downloads": downloads,
        "views": views,
    }
    return render(request, "resources/subject_dashboard.html", context)


# -------------------------------------------------------------------
# Report resource
# -------------------------------------------------------------------
@login_required
def report_resource(request, pk):
    resource = get_object_or_404(Resource, pk=pk)

    if request.method == "POST":
        form = ReportForm(request.POST)
        if form.is_valid():
            report = form.save(commit=False)
            report.resource = resource
            report.reporter = request.user
            report.save()

            if resource.owner != request.user:
                Notification.objects.create(
                    user=resource.owner,
                    notif_type="REPORT",
                    message=(
                        f"Your resource '{resource.title}' was reported by "
                        f"{request.user.username}."
                    ),
                    resource=resource,
                    report=report,
                )

            messages.success(request, "Thank you. Your report has been submitted.")
        else:
            messages.error(request, "Could not submit report, please check the form.")
    return redirect("resource_detail", pk=pk)


# -------------------------------------------------------------------
# My activity dashboard
# -------------------------------------------------------------------
@login_required
def my_activity(request):
    """
    Shows stats for the current user:
    - total uploads
    - favorites saved
    - ratings received
    - total views & downloads on their uploads
    - chart for views/downloads per upload
    - detailed table with each uploaded resource
    """
    user = request.user

    uploads_qs = Resource.objects.filter(owner=user).select_related("subject")

    uploads_count = uploads_qs.count()
    favorites_count = Favorite.objects.filter(user=user).count()
    ratings_received = Rating.objects.filter(resource__owner=user).count()

    total_views = uploads_qs.aggregate(total=Sum("view_count"))["total"] or 0
    total_downloads = uploads_qs.aggregate(total=Sum("download_count"))["total"] or 0

    # latest 10 uploads in chronological order (oldest first among those 10)
    uploads_for_chart = uploads_qs.order_by("-created_at")[:10][::-1]

    labels = [r.title[:20] for r in uploads_for_chart]
    views_data = [r.view_count for r in uploads_for_chart]
    downloads_data = [r.download_count for r in uploads_for_chart]

    my_uploads = uploads_qs.order_by("-created_at")

    context = {
        "uploads_count": uploads_count,
        "favorites_count": favorites_count,
        "ratings_received": ratings_received,
        "total_views": total_views,
        "total_downloads": total_downloads,
        "labels_json": json.dumps(labels),
        "views_json": json.dumps(views_data),
        "downloads_json": json.dumps(downloads_data),
        "my_uploads": my_uploads,
    }
    return render(request, "resources/my_activity.html", context)


# -------------------------------------------------------------------
# Admin analytics dashboard
# -------------------------------------------------------------------
@staff_member_required
def admin_analytics_dashboard(request):
    """
    Analytics for staff: uploads, downloads, reports, visits, active users.
    """
    # --- Totals ---
    total_resources = Resource.objects.count()
    total_downloads = Resource.objects.aggregate(total=Sum("download_count"))["total"] or 0
    total_views = Resource.objects.aggregate(total=Sum("view_count"))["total"] or 0
    total_reports = Report.objects.count()
    total_visits = Visit.objects.count()

    # --- Top resources by downloads ---
    top_resources = Resource.objects.order_by("-download_count")[:5]

    # --- Subject-wise uploads ---
    subject_stats = (
        Resource.objects.values("subject__name")
        .annotate(count=Count("id"))
        .order_by("-count")
    )

    # --- Uploads per day ---
    uploads_by_day_qs = (
        Resource.objects
        .annotate(day=TruncDate("created_at"))
        .values("day")
        .annotate(count=Count("id"))
        .order_by("day")
    )
    uploads_labels = [row["day"].strftime("%Y-%m-%d") for row in uploads_by_day_qs]
    uploads_counts = [row["count"] for row in uploads_by_day_qs]

    # --- Visits per day ---
    visits_by_day_qs = (
        Visit.objects
        .annotate(day=TruncDate("created_at"))
        .values("day")
        .annotate(count=Count("id"))
        .order_by("day")
    )
    visits_labels = [row["day"].strftime("%Y-%m-%d") for row in visits_by_day_qs]
    visits_counts = [row["count"] for row in visits_by_day_qs]

    # --- Downloads per subject ---
    subject_download_qs = (
        Resource.objects.values("subject__name")
        .annotate(downloads=Sum("download_count"))
        .order_by("-downloads")
    )
    subject_download_labels = [
        row["subject__name"] or "N/A" for row in subject_download_qs
    ]
    subject_download_counts = [row["downloads"] or 0 for row in subject_download_qs]

    # --- Most active users (by uploads) ---
    active_users_qs = (
        User.objects.annotate(
            uploads_count=Count("resources", distinct=True),
            ratings_given=Count("ratings", distinct=True),
            comments_made=Count("comments", distinct=True),
        )
        .order_by("-uploads_count")[:5]
    )

    context = {
        "total_resources": total_resources,
        "total_downloads": total_downloads,
        "total_views": total_views,
        "total_reports": total_reports,
        "total_visits": total_visits,
        "top_resources": top_resources,
        "subject_stats": subject_stats,
        "active_users": active_users_qs,
        "uploads_labels_json": json.dumps(uploads_labels),
        "uploads_counts_json": json.dumps(uploads_counts),
        "visits_labels_json": json.dumps(visits_labels),
        "visits_counts_json": json.dumps(visits_counts),
        "subject_download_labels_json": json.dumps(subject_download_labels),
        "subject_download_counts_json": json.dumps(subject_download_counts),
    }
    return render(request, "resources/admin_analytics.html", context)


# -------------------------------------------------------------------
# Notifications
# -------------------------------------------------------------------
@login_required
def notifications_list(request):
    notifications = request.user.notifications.order_by("-created_at")
    return render(
        request,
        "resources/notifications.html",
        {"notifications": notifications},
    )


@login_required
def notification_mark_read(request, pk):
    notif = get_object_or_404(Notification, pk=pk, user=request.user)
    notif.is_read = True
    notif.save()

    if notif.resource_id:
        return redirect("resource_detail", pk=notif.resource_id)
    return redirect("notifications_list")


# -------------------------------------------------------------------
# Verify resource (staff)
# -------------------------------------------------------------------
@staff_member_required
def verify_resource(request, pk):
    resource = get_object_or_404(Resource, pk=pk)

    if request.method == "POST":
        action = request.POST.get("action")
        note = request.POST.get("note", "").strip()

        if action not in ("APPROVED", "REJECTED"):
            messages.error(request, "Invalid action.")
            return redirect("resource_detail", pk=pk)

        resource.verification_status = action
        resource.verified_by = request.user
        resource.verified_at = timezone.now()
        resource.verification_note = note
        resource.save()

        msg = f"Your resource '{resource.title}' was "
        msg += "approved ✅" if action == "APPROVED" else "rejected ❌"
        if note:
            msg += f" – Note: {note}"

        Notification.objects.create(
            user=resource.owner,
            notif_type="REPORT_STATUS",
            message=msg,
            resource=resource,
        )

        messages.success(request, "Verification status updated.")
        return redirect("resource_detail", pk=pk)

    return redirect("resource_detail", pk=pk)


# -------------------------------------------------------------------
# Leaderboard
# -------------------------------------------------------------------
@login_required
def leaderboard(request):
    """
    Simple leaderboard for top contributors.
    Score = uploads*3 + total_downloads + comments_made
    """
    users = (
        User.objects.annotate(
            uploads_count=Count("resources", distinct=True),
            total_downloads=Sum("resources__download_count"),
            total_views=Sum("resources__view_count"),
            comments_made=Count("comments", distinct=True),
            ratings_given=Count("ratings", distinct=True),
        )
        .annotate(
            score=(
                F("uploads_count") * 3
                + Coalesce(F("total_downloads"), 0)
                + F("comments_made")
            )
        )
        .order_by("-score")[:20]
    )

    return render(request, "resources/leaderboard.html", {"users": users})
